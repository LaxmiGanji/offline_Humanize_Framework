"""
Enhanced file handler for processing different document formats
Supports: .txt, .pdf, .docx, .doc, .rtf, .xlsx, .csv, .pptx
With OCR support for scanned PDFs
"""

import os
import re
import subprocess
import sys

# Base required libraries
try:
    import PyPDF2
except ImportError:
    PyPDF2 = None

try:
    import docx
except ImportError:
    docx = None

# Optional libraries with graceful fallbacks
try:
    import pdfplumber
except ImportError:
    pdfplumber = None

try:
    import pandas as pd
except ImportError:
    pd = None

try:
    import csv
except ImportError:
    csv = None

try:
    from pptx import Presentation
except ImportError:
    Presentation = None

try:
    import pytesseract
    from PIL import Image
except ImportError:
    pytesseract = None
    Image = None

try:
    import fitz  # PyMuPDF
except ImportError:
    fitz = None

class FileHandler:
    """Handle different file types and extract text"""
    
    @staticmethod
    def read_file(file_path):
        """
        Read text from various file formats
        
        Args:
            file_path: Path to the file
        
        Returns:
            Extracted text content
        """
        if not os.path.exists(file_path):
            raise FileNotFoundError(f"File not found: {file_path}")
        
        file_extension = os.path.splitext(file_path)[1].lower()
        
        # Text files
        if file_extension == '.txt':
            return FileHandler._read_txt(file_path)
        
        # PDF files
        elif file_extension == '.pdf':
            return FileHandler._read_pdf_enhanced(file_path)
        
        # Word documents
        elif file_extension in ['.docx', '.doc']:
            return FileHandler._read_docx_enhanced(file_path)
        
        # Excel files
        elif file_extension in ['.xlsx', '.xls', '.csv']:
            return FileHandler._read_spreadsheet(file_path)
        
        # PowerPoint files
        elif file_extension in ['.pptx', '.ppt']:
            return FileHandler._read_powerpoint(file_path)
        
        # Rich text format
        elif file_extension == '.rtf':
            return FileHandler._read_rtf(file_path)
        
        else:
            return f"Unsupported file format: {file_extension}. Supported formats: .txt, .pdf, .docx, .xlsx, .csv, .pptx"
    
    @staticmethod
    def _read_txt(file_path):
        """Read text file"""
        try:
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as file:
                return file.read()
        except Exception as e:
            return f"Error reading text file: {str(e)}"
    
    @staticmethod
    def _read_pdf_enhanced(file_path):
        """Enhanced PDF reading with multiple methods"""
        text = []
        error_messages = []
        
        # Method 1: Try PyPDF2 (always available if installed)
        if PyPDF2:
            try:
                with open(file_path, 'rb') as file:
                    pdf_reader = PyPDF2.PdfReader(file)
                    for page in pdf_reader.pages:
                        page_text = page.extract_text()
                        if page_text and page_text.strip():
                            text.append(page_text)
                if text:
                    return '\n\n'.join(text)
            except Exception as e:
                error_messages.append(f"PyPDF2 failed: {e}")
        
        # Method 2: Try pdfplumber (if available)
        if pdfplumber and not text:
            try:
                with pdfplumber.open(file_path) as pdf:
                    for page in pdf.pages:
                        page_text = page.extract_text()
                        if page_text and page_text.strip():
                            text.append(page_text)
                if text:
                    return '\n\n'.join(text)
            except Exception as e:
                error_messages.append(f"pdfplumber failed: {e}")
        
        # Method 3: Try PyMuPDF (if available)
        if fitz and not text:
            try:
                doc = fitz.open(file_path)
                for page in doc:
                    page_text = page.get_text()
                    if page_text and page_text.strip():
                        text.append(page_text)
                if text:
                    return '\n\n'.join(text)
            except Exception as e:
                error_messages.append(f"PyMuPDF failed: {e}")
        
        # Method 4: Try OCR (if available)
        if pytesseract and Image and fitz and not text:
            try:
                doc = fitz.open(file_path)
                for page_num in range(len(doc)):
                    page = doc[page_num]
                    # Convert page to image
                    pix = page.get_pixmap()
                    img_data = pix.tobytes("png")
                    img = Image.open(io.BytesIO(img_data))
                    
                    # Perform OCR
                    ocr_text = pytesseract.image_to_string(img)
                    if ocr_text and ocr_text.strip():
                        text.append(f"[Page {page_num + 1} - OCR Extracted]\n{ocr_text}")
                if text:
                    return '\n\n'.join(text)
            except Exception as e:
                error_messages.append(f"OCR failed: {e}")
        
        if not text:
            error_msg = "No text could be extracted from the PDF. "
            if error_messages:
                error_msg += "Errors: " + "; ".join(error_messages)
            else:
                error_msg += "The PDF might be scanned or image-based. Install additional libraries: pip install pdfplumber pymupdf pytesseract pillow"
            return error_msg
        
        return '\n\n'.join(text)
    
    @staticmethod
    def _read_docx_enhanced(file_path):
        """Enhanced Word document reading"""
        if not docx:
            return "python-docx library not installed. Please install: pip install python-docx"
        
        text = []
        
        try:
            doc = docx.Document(file_path)
            
            # Extract from paragraphs
            for paragraph in doc.paragraphs:
                if paragraph.text:
                    text.append(paragraph.text)
            
            # Extract from tables
            for table in doc.tables:
                for row in table.rows:
                    row_text = []
                    for cell in row.cells:
                        if cell.text:
                            row_text.append(cell.text)
                    if row_text:
                        text.append(' | '.join(row_text))
            
        except Exception as e:
            return f"Error reading DOCX: {str(e)}"
        
        if not text:
            return "No text could be extracted from the DOCX file."
        
        return '\n'.join(text)
    
    @staticmethod
    def _read_spreadsheet(file_path):
        """Read Excel/CSV files"""
        text = []
        file_extension = os.path.splitext(file_path)[1].lower()
        
        try:
            if file_extension == '.csv':
                # Read CSV
                if csv:
                    with open(file_path, 'r', encoding='utf-8', errors='ignore') as file:
                        csv_reader = csv.reader(file)
                        for row in csv_reader:
                            if row:
                                text.append(' | '.join(str(cell) for cell in row if cell))
                else:
                    return "CSV module not available"
            
            elif pd is not None:  # Excel files with pandas
                df_dict = pd.read_excel(file_path, sheet_name=None)
                for sheet_name, df in df_dict.items():
                    text.append(f"\n--- Sheet: {sheet_name} ---\n")
                    # Convert dataframe to text
                    for _, row in df.iterrows():
                        row_text = ' | '.join(str(val) for val in row if pd.notna(val))
                        if row_text:
                            text.append(row_text)
            else:
                return "pandas library not installed. Please install: pip install pandas openpyxl"
        
        except Exception as e:
            return f"Error reading spreadsheet: {str(e)}"
        
        return '\n'.join(text)
    
    @staticmethod
    def _read_powerpoint(file_path):
        """Read PowerPoint files"""
        if not Presentation:
            return "python-pptx library not installed. Please install: pip install python-pptx"
        
        text = []
        
        try:
            prs = Presentation(file_path)
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_text = [f"--- Slide {slide_num} ---"]
                
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        slide_text.append(shape.text)
                    
                    # Check for tables
                    if hasattr(shape, "table") and shape.table:
                        for row in shape.table.rows:
                            row_text = []
                            for cell in row.cells:
                                if cell.text:
                                    row_text.append(cell.text)
                            if row_text:
                                slide_text.append(' | '.join(row_text))
                
                if len(slide_text) > 1:  # Only add if there's content
                    text.append('\n'.join(slide_text))
        
        except Exception as e:
            return f"Error reading PowerPoint: {str(e)}"
        
        return '\n'.join(text)
    
    @staticmethod
    def _read_rtf(file_path):
        """Read RTF files (simple version)"""
        try:
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as file:
                content = file.read()
                # Remove RTF formatting (very basic)
                content = re.sub(r'\\[a-z]+[\d-]*', ' ', content)
                content = re.sub(r'\{|\}', ' ', content)
                content = re.sub(r'\s+', ' ', content)
                return content
        except Exception as e:
            return f"Error reading RTF: {str(e)}"
    
    @staticmethod
    def get_file_info(file_path):
        """Get file information"""
        file_extension = os.path.splitext(file_path)[1].lower()
        file_size = os.path.getsize(file_path)
        
        # Convert size to readable format
        if file_size < 1024:
            size_str = f"{file_size} B"
        elif file_size < 1024 * 1024:
            size_str = f"{file_size / 1024:.1f} KB"
        else:
            size_str = f"{file_size / (1024 * 1024):.1f} MB"
        
        return {
            'name': os.path.basename(file_path),
            'extension': file_extension,
            'size': size_str,
            'size_bytes': file_size
        }
    
    @staticmethod
    def get_file_preview(file_path, max_chars=500):
        """Get preview of file content"""
        try:
            text = FileHandler.read_file(file_path)
            if len(text) > max_chars:
                return text[:max_chars] + "..."
            return text
        except:
            return "Preview not available"